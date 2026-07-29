import streamlit as st
import gspread
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
from datetime import datetime
from io import BytesIO

# --- dependensi ekspor PPTX (pip install python-pptx) ---
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

st.set_page_config(page_title="Dashboard Overview Project", layout="wide")

# ═══════════════════════════════════════════════════════════
# KONFIGURASI
# ═══════════════════════════════════════════════════════════
ID_SPREADSHEET_TRACKER = '198lmRiSC1VjZeEL9CK11miLLQlBMQBjXKubEJMX0MO0'

# Definisi tiap project: sheet tracker + sheet issue log + cara baca
# "progress_mode" menentukan bagaimana % penyelesaian dihitung:
#   - "status"  : hitung baris Done vs total baris aktif
#   - "numeric" : rata-rata kolom progres numerik (0..1)
PROJECTS = {
    "DSSP": {
        "tracker_sheet": "Tracker + Timeline DSSP",
        "issue_sheet":   "ISSUE LOG DSSP",
        "progress_mode": "status",
        "extra_selesai": ("under review",),  # Under Review dihitung sebagai selesai
        "warna": "#3498db",
    },
    "KCS CIMB Batam": {
        "tracker_sheet": "Tracker KCS Batam",
        "issue_sheet":   "ISSUE LOG CIMB Batam",
        "progress_mode": "numeric",
        "warna": "#9b59b6",
    },
    "KDS": {
        "tracker_sheet": "Tracker KDS",
        "issue_sheet":   "ISSUE LOG KDS",
        "progress_mode": "status",
        "warna": "#e67e22",
    },
    "Nova": {
        "tracker_sheet": "Tracker Nova",
        "issue_sheet":   "ISSUE LOG NOVA",
        "progress_mode": "status",
        "warna": "#1abc9c",
    },
    "Uniqlo": {
        "tracker_sheet": "Tracker VOIP Uniqlo",
        "issue_sheet":   "ISSUE LOG Uniqlo",
        "progress_mode": "status",
        "warna": "#e74c3c",
    },
}

# Nilai status yang dianggap "selesai" (lower-case matching)
STATUS_SELESAI = ('done', 'selesai', 'closed', 'invoiced', 'close', 'complete')
# Nilai status yang dianggap "belum mulai / tidak dihitung sebagai aktif"
STATUS_BELUM = ('tbd', 'nys', 'not yet', 'not started', 'belum', 'pending')


# ═══════════════════════════════════════════════════════════
# KONEKSI
# ═══════════════════════════════════════════════════════════
def _get_client():
    credentials = dict(st.secrets["gcp_service_account"])
    return gspread.service_account_from_dict(credentials)


def _find_header_row(raw, keywords, max_scan=15):
    """Cari index baris header berdasarkan skor kata kunci."""
    best_idx, best_score = 0, -1
    for i, row in enumerate(raw[:max_scan]):
        filled = [str(c).strip() for c in row if str(c).strip()]
        kw_bonus = sum(1 for c in filled if any(k in str(c).lower() for k in keywords))
        score = len(filled) + kw_bonus * 2
        if score > best_score:
            best_score, best_idx = score, i
    return best_idx


def _dedup_columns(df):
    df = df.loc[:, [bool(str(c).strip()) for c in df.columns]]
    df = df.loc[:, ~pd.Index(df.columns).duplicated()]
    return df


def _cari_kolom(cols, *keywords):
    # Prioritas 1: exact match (case-insensitive) — hindari 'Status' ketemu 'Remark Status'
    for kw in keywords:
        for c in cols:
            if str(c).strip().lower() == kw.lower():
                return c
    # Prioritas 2: substring match
    for kw in keywords:
        for c in cols:
            if kw.lower() in str(c).lower():
                return c
    return None


def parse_rupiah(series):
    """Konversi kolom teks nominal (mis. 'Rp 500.000', '500000', '45,000') ke numeric."""
    return pd.to_numeric(
        series.astype(str)
              .str.replace(r'[Rp\s\.]', '', regex=True)
              .str.replace(',', '', regex=False),
        errors='coerce'
    ).fillna(0)


def compute_margin(df, nama_project):
    """Hitung rincian revenue/cost/margin per baris untuk KDS & Nova.
    Kembalikan (df_rincian, total_rev, total_cost, total_margin) atau (None, ...) bila N/A.
    - KDS : pakai kolom Revenue, Cost, Margin yg sudah ada.
    - Nova: Revenue = PO, Cost = Actual Cost, Margin = PO - Actual Cost.
    """
    if df.empty:
        return None, 0, 0, 0
    cols = df.columns.tolist()

    if 'KDS' in nama_project:
        col_rev  = _cari_kolom(cols, 'revenue')
        col_cost = next((c for c in cols if c.strip().lower() == 'cost'), None) or _cari_kolom(cols, 'cost')
        col_margin = next((c for c in cols if c.strip().lower() == 'margin'), None)
        col_label = _cari_kolom(cols, 'kerjaan', 'customer', 'assignment')
        if not (col_rev and col_cost):
            return None, 0, 0, 0
        rev = parse_rupiah(df[col_rev])
        cost = parse_rupiah(df[col_cost])
        margin = parse_rupiah(df[col_margin]) if col_margin else (rev - cost)
        label = df[col_label].astype(str) if col_label else pd.Series(range(len(df))).astype(str)

    elif 'Nova' in nama_project:
        col_po = next((c for c in cols if c.strip().lower() == 'po'), None) or _cari_kolom(cols, 'po')
        col_ac = _cari_kolom(cols, 'actual cost')
        col_label = _cari_kolom(cols, 'nama entitas', 'nama daerah', 'entitas', 'daerah')
        if not (col_po and col_ac):
            return None, 0, 0, 0
        rev = parse_rupiah(df[col_po])
        cost = parse_rupiah(df[col_ac])
        margin = rev - cost
        label = df[col_label].astype(str) if col_label else pd.Series(range(len(df))).astype(str)

    else:
        return None, 0, 0, 0

    rincian = pd.DataFrame({
        'Lokasi': label.values,
        'Revenue': rev.values,
        'Cost': cost.values,
        'Margin': margin.values,
    })
    # buang baris yg semua nilainya nol (baris kosong padding)
    rincian = rincian[(rincian['Revenue'] != 0) | (rincian['Cost'] != 0) | (rincian['Margin'] != 0)]
    rincian['Margin %'] = rincian.apply(
        lambda r: (r['Margin'] / r['Revenue'] * 100) if r['Revenue'] else 0, axis=1)
    return (rincian.reset_index(drop=True),
            float(rincian['Revenue'].sum()),
            float(rincian['Cost'].sum()),
            float(rincian['Margin'].sum()))


def compute_cost_health(df, nama_project):
    """Analisa kepatuhan anggaran (Plan vs Actual Cost) untuk cross-subsidy.

    Berbeda dari compute_margin (yg bicara revenue vs cost / profitabilitas).
    Di sini fokusnya: apakah actual cost masih di dalam plan, dan apakah
    lokasi yg overcost masih bisa ditutup oleh sisa anggaran (headroom) dari
    lokasi yg undercost.

    Return dict:
      realized     : DataFrame lokasi yg sudah punya Actual Cost (>0)
      unrealized   : DataFrame lokasi yg Actual Cost masih kosong
      plan_real    : total plan dari lokasi terealisasi
      actual_real  : total actual dari lokasi terealisasi
      headroom     : plan_real - actual_real (positif = masih ada sisa pool)
      over_total   : total nilai kelebihan biaya dari lokasi overcost saja
      n_over       : jumlah lokasi overcost
    Atau None bila kolom tak tersedia.
    """
    if df.empty:
        return None
    cols = df.columns.tolist()
    col_plan   = _cari_kolom(cols, 'plan cost')
    col_actual = _cari_kolom(cols, 'actual cost')
    col_label  = _cari_kolom(cols, 'store', 'lokasi', 'site', 'nama')
    col_city   = _cari_kolom(cols, 'city', 'kota', 'daerah')
    if not (col_plan and col_actual):
        return None

    plan   = parse_rupiah(df[col_plan])
    actual = parse_rupiah(df[col_actual])
    label  = (df[col_label].astype(str) if col_label
              else pd.Series(range(len(df))).astype(str))
    city   = df[col_city].astype(str) if col_city else pd.Series([''] * len(df))

    base = pd.DataFrame({
        'Lokasi': label.values,
        'Kota':   city.values,
        'Plan':   plan.values,
        'Actual': actual.values,
    })
    # buang baris padding kosong (tanpa nama & tanpa plan)
    base = base[(base['Lokasi'].str.strip() != '') & (base['Plan'] != 0)]

    # Lokasi "terealisasi" = actual sudah diisi (>0). Actual==0 dianggap belum,
    # sesuai keputusan: lokasi belum realisasi ditampilkan terpisah.
    realized   = base[base['Actual'] > 0].copy()
    unrealized = base[base['Actual'] <= 0].copy()

    if not realized.empty:
        realized['Variance'] = realized['Actual'] - realized['Plan']   # + = overcost
        realized['Variance %'] = realized.apply(
            lambda r: (r['Variance'] / r['Plan'] * 100) if r['Plan'] else 0, axis=1)
        realized['Kondisi'] = realized['Variance'].apply(
            lambda v: 'Overcost' if v > 0 else ('Undercost' if v < 0 else 'Tepat'))

    plan_real   = float(realized['Plan'].sum())   if not realized.empty else 0.0
    actual_real = float(realized['Actual'].sum()) if not realized.empty else 0.0
    headroom    = plan_real - actual_real
    over_total  = float(realized.loc[realized['Variance'] > 0, 'Variance'].sum()) \
                  if not realized.empty else 0.0
    n_over      = int((realized['Variance'] > 0).sum()) if not realized.empty else 0

    return {
        'realized': realized.reset_index(drop=True),
        'unrealized': unrealized.reset_index(drop=True),
        'plan_real': plan_real,
        'actual_real': actual_real,
        'headroom': headroom,
        'over_total': over_total,
        'n_over': n_over,
    }


@st.cache_data(ttl=60)          # FIX: sebelumnya tanpa decorator, jadi
def load_tracker(nama_sheet: str) -> pd.DataFrame:   # load_tracker.clear() error
    gc = _get_client()
    sh = gc.open_by_key(ID_SPREADSHEET_TRACKER)
    try:
        ws = sh.worksheet(nama_sheet)
    except gspread.exceptions.WorksheetNotFound:
        return pd.DataFrame()

    raw = ws.get_all_values()
    if not raw:
        return pd.DataFrame()

    KATA_KUNCI = (
        'no', 'nama', 'status', 'tanggal', 'kode', 'site', 'kategori', 'po',
        'plan cost', 'actual cost', 'progres', 'progress', 'target', 'realisasi',
        'harga', 'cost', 'persentase', 'revenue', 'margin', 'period', 'remark',
        'deadline', 'prioritas', 'customer', 'spk', 'pekerjaan', 'tagihan', 'id',
        'pic', 'assign', 'fase', 'deskripsi', 'sistem', 'catatan', 'checklist',
        'task', 'sub-task', 'entitas', 'daerah'
    )
    # Anchor eksplisit: baris header sebenarnya adalah baris yang memuat 'Task' + 'Sub-Task'
    # (khusus sheet DSSP yg punya sub-header baris day-number yg membingungkan skor).
    hdr = None
    for i, row in enumerate(raw[:15]):
        low = [str(c).strip().lower() for c in row]
        if 'sub-task' in low and 'task' in low:
            hdr = i
            break
    if hdr is None:
        hdr = _find_header_row(raw, KATA_KUNCI)

    # Kasus khusus DSSP: header terpecah 2 baris. Kolom 'Status' asli ada di
    # baris hdr+1 (di bawah grup 'Delay'); baris hdr hanya punya 'Remark Status'.
    # Gabungkan kedua baris agar kolom 'Status' tertangkap dengan benar.
    next_row = raw[hdr + 1] if hdr + 1 < len(raw) else []
    is_split_header = any(str(c).strip().lower() == 'status' for c in next_row)

    if is_split_header:
        h1 = [str(c).strip() for c in raw[hdr]]
        h2 = [str(c).strip() for c in next_row]
        merged = [b if b else a for a, b in zip(h1, h2)]
        rows = []
        for r in raw[hdr + 2:]:
            if not any(str(r[j]).strip() for j in range(1, min(13, len(r)))):
                break
            rows.append(r)
        df = pd.DataFrame(rows, columns=merged)
    else:
        headers = [str(h).strip() for h in raw[hdr]]
        df = pd.DataFrame(raw[hdr + 1:], columns=headers)

    df = _dedup_columns(df)
    df = df[df.apply(lambda r: any(str(v).strip() for v in r), axis=1)]
    return df.reset_index(drop=True)


@st.cache_data(ttl=60)
def load_issue_log(nama_sheet: str) -> pd.DataFrame:
    gc = _get_client()
    sh = gc.open_by_key(ID_SPREADSHEET_TRACKER)
    try:
        ws = sh.worksheet(nama_sheet)
    except gspread.exceptions.WorksheetNotFound:
        return pd.DataFrame()

    raw = ws.get_all_values()
    if not raw:
        return pd.DataFrame()

    hdr = _find_header_row(raw, ('no', 'deskripsi', 'status', 'pic', 'remark', 'tanggal', 'ticket'))
    headers = [str(h).strip() for h in raw[hdr]]
    df = pd.DataFrame(raw[hdr + 1:], columns=headers)
    df = _dedup_columns(df)

    # Baris issue valid = punya deskripsi (bukan sekadar nomor urut kosong)
    col_desc = _cari_kolom(df.columns, 'deskripsi', 'description', 'issue')
    if col_desc:
        df = df[df[col_desc].astype(str).str.strip() != '']
    else:
        df = df[df.apply(lambda r: sum(1 for v in r if str(v).strip()) >= 2, axis=1)]
    return df.reset_index(drop=True)


# ═══════════════════════════════════════════════════════════
# ANALISA: hitung progress & issue untuk tiap project
# ═══════════════════════════════════════════════════════════
def hitung_progress(df: pd.DataFrame, mode: str, extra_selesai=()):
    """Kembalikan (persen, total_aktif, jumlah_selesai, breakdown_status_dict).
    extra_selesai: status tambahan (lower-case) yg dianggap selesai utk project ini,
    mis. ('under review',) untuk DSSP."""
    if df.empty:
        return 0.0, 0, 0, {}

    cols = df.columns.tolist()

    if mode == "numeric":
        col_prog = _cari_kolom(cols, 'progres', 'progress')
        col_status = _cari_kolom(cols, 'status')
        if col_prog:
            prog = pd.to_numeric(
                df[col_prog].astype(str).str.replace('%', '', regex=False).str.strip(),
                errors='coerce'
            )
            prog = prog.dropna()
            # Jika nilai > 1 anggap skala 0-100, normalisasi ke 0-1
            if not prog.empty and prog.max() > 1.5:
                prog = prog / 100.0
            persen = float(prog.mean() * 100) if not prog.empty else 0.0
            total = len(prog)
            selesai = int((prog >= 0.999).sum())
            breakdown = {}
            if col_status:
                breakdown = (df[col_status].replace(r'^\s*$', 'Unknown', regex=True)
                             .fillna('Unknown').value_counts().to_dict())
            return persen, total, selesai, breakdown

    # mode == "status" (default)
    col_status = (_cari_kolom(cols, 'status pekerjaan')
                  or _cari_kolom(cols, 'status'))
    if not col_status:
        return 0.0, len(df), 0, {}

    s = df[col_status].astype(str).str.strip()
    s = s[s != '']
    # Baris "belum mulai/TBD/NYS" tetap dihitung sebagai bagian dari scope aktif
    total = len(s)
    if total == 0:
        return 0.0, 0, 0, {}

    kriteria = tuple(STATUS_SELESAI) + tuple(extra_selesai)
    selesai = int(s.str.lower().apply(lambda x: any(k in x for k in kriteria)).sum())
    persen = selesai / total * 100
    breakdown = s.value_counts().to_dict()
    return persen, total, selesai, breakdown


def hitung_issue(df: pd.DataFrame):
    """Kembalikan (total, open, close)."""
    if df.empty:
        return 0, 0, 0
    col_status = _cari_kolom(df.columns, 'status')
    total = len(df)
    if not col_status:
        return total, 0, 0
    s = df[col_status].astype(str).str.strip().str.lower()
    close = int(s.apply(lambda x: 'close' in x).sum())
    open_ = int(s.apply(lambda x: x == 'open' or 'open' in x).sum())
    # sisa yg tidak jelas -> tidak dihitung open/close
    return total, open_, close


# ═══════════════════════════════════════════════════════════
# EKSPOR PPTX
# ═══════════════════════════════════════════════════════════
# Chart dibuat NATIVE PowerPoint (bukan gambar Plotly), jadi tidak butuh
# kaleido / headless Chrome dan hasilnya bisa diedit langsung di PowerPoint.

INK      = RGBColor(0x1B, 0x24, 0x30)   # teks utama / background cover
INK_SOFT = RGBColor(0x44, 0x50, 0x5E)
MUTED    = RGBColor(0x76, 0x82, 0x8F)
LINE     = RGBColor(0xDD, 0xE2, 0xE7)
PAPER    = RGBColor(0xFF, 0xFF, 0xFF)
BAND     = RGBColor(0xF4, 0xF6, 0xF8)
OK       = RGBColor(0x2E, 0x9E, 0x5B)
WARN     = RGBColor(0xD9, 0x82, 0x2B)
BAD      = RGBColor(0xC0, 0x39, 0x2B)

MAX_ROW_PER_SLIDE = 7    # baris tabel isu per slide
DESC_MAX = 95            # potong deskripsi isu biar tidak overflow

BULAN_ID = ["Januari", "Februari", "Maret", "April", "Mei", "Juni", "Juli",
            "Agustus", "September", "Oktober", "November", "Desember"]


def _hex_to_rgb(h):
    h = str(h).lstrip("#")
    try:
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    except Exception:
        return MUTED


def _potong(teks, n=DESC_MAX):
    t = " ".join(str(teks).split())
    return t if len(t) <= n else t[: n - 1].rstrip() + "…"


def _tanggal(v):
    """Normalkan nilai tanggal (Sheets sering mengirim '2026-07-23 00:00:00')."""
    t = str(v).strip()
    if not t or t.lower() in ("nan", "nat", "none"):
        return "—"
    try:
        # '2026-07-23...' = ISO (bulan di tengah); '23/07/2026' = hari di depan
        dayfirst = not t[:4].isdigit()
        return pd.to_datetime(t, errors="raise", dayfirst=dayfirst).strftime("%d %b %Y")
    except Exception:
        return _potong(t, 12)


def _txt(slide, x, y, w, h, teks, size=14, bold=False, color=INK,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if spacing:
        p.line_spacing = spacing
    r = p.add_run()
    r.text = str(teks)
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = "Calibri"
    return box


def _slide_kosong(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])   # blank layout


def _header(slide, kicker, judul, no_halaman=None):
    """Header standar: kicker kecil di atas, judul besar."""
    _txt(slide, 0.7, 0.45, 11.9, 0.28, kicker.upper(), size=11, bold=True, color=MUTED)
    _txt(slide, 0.7, 0.78, 11.9, 0.55, judul, size=26, bold=True, color=INK)
    if no_halaman is not None:
        _txt(slide, 12.0, 6.95, 0.6, 0.3, str(no_halaman), size=10,
             color=MUTED, align=PP_ALIGN.RIGHT)


def _kpi(slide, x, y, w, angka, label, warna=INK):
    _txt(slide, x, y, w, 0.75, angka, size=40, bold=True, color=warna)
    _txt(slide, x, y + 0.72, w, 0.55, label, size=11, color=MUTED)


def _tabel(slide, x, y, w, header, rows, col_w=None, font=10.5,
           warna_kolom_status=None):
    """Tabel rapi tanpa style bawaan PowerPoint."""
    n_r, n_c = len(rows) + 1, len(header)
    h = Inches(0.38 + 0.34 * len(rows))
    shp = slide.shapes.add_table(n_r, n_c, Inches(x), Inches(y), Inches(w), h)
    tbl = shp.table
    tbl.first_row = True
    tbl.horz_banding = False

    if col_w:
        total = sum(col_w)
        for i, cw in enumerate(col_w):
            tbl.columns[i].width = Emu(int(Inches(w) * cw / total))

    tbl.rows[0].height = Inches(0.38)
    for j, teks in enumerate(header):
        c = tbl.cell(0, j)
        c.text = str(teks)
        c.fill.solid()
        c.fill.fore_color.rgb = INK
        c.margin_left = c.margin_right = Inches(0.08)
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = c.text_frame.paragraphs[0]
        p.font.size = Pt(font)
        p.font.bold = True
        p.font.color.rgb = PAPER
        p.font.name = "Calibri"

    for i, row in enumerate(rows, start=1):
        tbl.rows[i].height = Inches(0.34)
        for j, teks in enumerate(row):
            c = tbl.cell(i, j)
            c.text = str(teks)
            c.fill.solid()
            c.fill.fore_color.rgb = PAPER if i % 2 else BAND
            c.margin_left = c.margin_right = Inches(0.08)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = c.text_frame.paragraphs[0]
            p.font.size = Pt(font)
            p.font.name = "Calibri"
            p.font.color.rgb = INK_SOFT
            if warna_kolom_status is not None and j == warna_kolom_status:
                p.font.bold = True
                p.font.color.rgb = BAD if "open" in str(teks).lower() else OK
    return shp


def _rapikan_chart(chart, label_size=10):
    chart.font.size = Pt(10)
    chart.font.name = "Calibri"
    chart.font.color.rgb = INK_SOFT
    chart.has_title = False
    try:
        chart.value_axis.has_major_gridlines = True
        chart.value_axis.major_gridlines.format.line.color.rgb = LINE
        chart.value_axis.major_gridlines.format.line.width = Pt(0.5)
        chart.category_axis.has_major_gridlines = False
        chart.value_axis.tick_labels.font.size = Pt(label_size)
        chart.category_axis.tick_labels.font.size = Pt(label_size)
    except Exception:
        pass


def _slide_cover(prs, data, judul, tanggal):
    s = _slide_kosong(prs)
    bg = s.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = INK

    total_issue = sum(d["iss_total"] for d in data.values())
    total_open = sum(d["iss_open"] for d in data.values())
    avg = sum(d["persen"] for d in data.values()) / len(data) if data else 0

    _txt(s, 0.9, 1.05, 11.5, 0.3, "PROGRESS & ISSUE REPORT", size=12,
         bold=True, color=RGBColor(0x8E, 0x9B, 0xA8))
    _txt(s, 0.9, 1.55, 9.5, 1.9, judul, size=44, bold=True, color=PAPER, spacing=1.05)
    _txt(s, 0.9, 3.55, 9.5, 0.35,
         f"{len(data)} proyek dalam pengawalan   |   Data per {tanggal}",
         size=13, color=RGBColor(0xA8, 0xB4, 0xC0))

    kpis = [
        (f"{avg:.0f}%", "rata-rata progres"),
        (str(total_open), "isu terbuka"),
        (str(total_issue), "total isu tercatat"),
        (str(len(data)), "proyek aktif"),
    ]
    for i, (angka, label) in enumerate(kpis):
        x = 0.9 + i * 2.95
        _txt(s, x, 4.5, 2.7, 0.8, angka, size=38, bold=True,
             color=BAD if label == "isu terbuka" and total_open else PAPER)
        _txt(s, x, 5.32, 2.7, 0.4, label, size=11, color=RGBColor(0xA8, 0xB4, 0xC0))

    _txt(s, 0.9, 6.6, 11.5, 0.6,
         "Sumber: Tracker Project & Issue Log (Google Sheets). "
         "Angka dihasilkan otomatis dari dashboard pada saat ekspor.",
         size=9.5, color=RGBColor(0x7A, 0x86, 0x93))
    return s


def _slide_ringkasan(prs, data, hal):
    s = _slide_kosong(prs)
    _header(s, "Ringkasan", "Progres dan isu terbuka seluruh proyek", hal)

    header = ["Proyek", "Progres", "Item selesai", "Total isu", "Isu terbuka"]
    rows = []
    for nama, d in data.items():
        rows.append([
            nama,
            f"{d['persen']:.0f}%",
            f"{d['selesai']}/{d['total_aktif']}",
            str(d["iss_total"]),
            str(d["iss_open"]) if d["iss_open"] else "—",
        ])
    _tabel(s, 0.7, 1.65, 11.9, header, rows, col_w=[3.4, 1.4, 1.7, 1.3, 1.5])

    y = 1.65 + 0.38 + 0.34 * len(rows) + 0.45
    perhatian = [n for n, d in data.items() if d["iss_open"] > 0]
    tertinggal = [n for n, d in data.items() if d["persen"] < 60]

    if perhatian:
        _txt(s, 0.7, y, 11.9, 0.32,
             f"Butuh tindak lanjut: {', '.join(perhatian)}", size=15, bold=True, color=BAD)
        y += 0.42
    if tertinggal:
        _txt(s, 0.7, y, 11.9, 0.32,
             f"Progres di bawah 60%: {', '.join(tertinggal)}", size=13, color=WARN)
        y += 0.38

    _txt(s, 0.7, 6.6, 11.9, 0.5,
         "Metode: progres berbasis status (rasio item selesai) kecuali KCS Batam "
         "yang memakai rata-rata kolom Progres numerik. Unit kerja tiap proyek "
         "berbeda — bandingkan tren, bukan angka absolut antarproyek.",
         size=9.5, color=MUTED)
    return s


def _slide_chart_progress(prs, data, hal):
    s = _slide_kosong(prs)
    _header(s, "Grafik", "Persentase penyelesaian per proyek", hal)

    nama = list(data.keys())
    cd = CategoryChartData()
    cd.categories = nama
    cd.add_series("Progres", [round(d["persen"], 1) for d in data.values()])

    gf = s.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.7),
                            Inches(1.7), Inches(11.9), Inches(4.5), cd)
    chart = gf.chart
    chart.has_legend = False
    _rapikan_chart(chart)
    plot = chart.plots[0]
    plot.gap_width = 60
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.number_format = '0"%"'
    dl.number_format_is_linked = False
    dl.font.size = Pt(11)
    dl.font.bold = True
    dl.font.color.rgb = INK
    try:
        dl.position = XL_LABEL_POSITION.OUTSIDE_END
    except Exception:
        pass
    for i, n in enumerate(nama):
        pt = plot.series[0].points[i]
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = _hex_to_rgb(data[n]["cfg"].get("warna", "#3498db"))
    chart.value_axis.maximum_scale = 100.0
    chart.value_axis.minimum_scale = 0.0

    _txt(s, 0.7, 6.5, 11.9, 0.4,
         "Nilai 100% berarti seluruh item pada tracker berstatus selesai, "
         "bukan berarti proyek sudah serah terima.", size=9.5, color=MUTED)
    return s


def _slide_chart_issue(prs, data, hal):
    s = _slide_kosong(prs)
    _header(s, "Grafik", "Status isu: terbuka vs tertutup", hal)

    nama = list(data.keys())
    cd = CategoryChartData()
    cd.categories = nama
    cd.add_series("Open", [d["iss_open"] for d in data.values()])
    cd.add_series("Close", [d["iss_close"] for d in data.values()])

    gf = s.shapes.add_chart(XL_CHART_TYPE.BAR_STACKED, Inches(0.7),
                            Inches(1.7), Inches(11.9), Inches(4.4), cd)
    chart = gf.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.TOP
    chart.legend.include_in_layout = False
    _rapikan_chart(chart)
    plot = chart.plots[0]
    plot.gap_width = 60
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.font.size = Pt(10)
    dl.font.color.rgb = PAPER
    dl.position = XL_LABEL_POSITION.CENTER   # stacked: HARUS ctr/inEnd/inBase
    plot.series[0].format.fill.solid()
    plot.series[0].format.fill.fore_color.rgb = BAD
    plot.series[1].format.fill.solid()
    plot.series[1].format.fill.fore_color.rgb = OK

    _txt(s, 0.7, 6.45, 11.9, 0.4,
         "Isu tanpa status yang jelas pada issue log tidak dihitung sebagai "
         "open maupun close.", size=9.5, color=MUTED)
    return s


def _kolom_issue(df):
    """Deteksi kolom issue log secara longgar (nama kolom antar-sheet berbeda)."""
    def cari(*kw):
        for k in kw:
            for c in df.columns:
                if str(c).strip().lower() == k:
                    return c
        for k in kw:
            for c in df.columns:
                if k in str(c).lower():
                    return c
        return None
    return {
        "status": cari("status"),
        "desc": cari("deskripsi", "description", "issue", "permasalahan"),
        "pic": cari("pic", "penanggung"),
        "tgl": cari("tanggal", "date"),
    }


def _slide_issue_open(prs, data, hal_mulai):
    """Satu (atau lebih) slide berisi daftar isu OPEN lintas proyek."""
    rows = []
    for nama, d in data.items():
        df = d.get("df_issue")
        if df is None or df.empty:
            continue
        k = _kolom_issue(df)
        if not k["status"]:
            continue
        mask = df[k["status"]].astype(str).str.strip().str.lower().str.contains('open', na=False)
        for _, r in df[mask].iterrows():
            rows.append([
                nama,
                _tanggal(r[k["tgl"]]) if k["tgl"] else "—",
                _potong(r[k["desc"]]) if k["desc"] else "—",
                _potong(r[k["pic"]], 18) if k["pic"] else "—",
            ])

    hal = hal_mulai
    if not rows:
        s = _slide_kosong(prs)
        _header(s, "Isu terbuka", "Tidak ada isu terbuka", hal)
        _txt(s, 0.7, 2.6, 11.9, 0.6,
             "Seluruh isu yang tercatat pada issue log sudah berstatus close.",
             size=18, bold=True, color=OK)
        return hal + 1

    total = len(rows)
    chunks = [rows[i:i + MAX_ROW_PER_SLIDE] for i in range(0, total, MAX_ROW_PER_SLIDE)]
    for idx, chunk in enumerate(chunks, start=1):
        s = _slide_kosong(prs)
        sub = f" ({idx}/{len(chunks)})" if len(chunks) > 1 else ""
        _header(s, "Isu terbuka", f"{total} isu menunggu tindak lanjut{sub}", hal)
        _tabel(s, 0.7, 1.7, 11.9,
               ["Proyek", "Tanggal", "Deskripsi isu", "PIC"], chunk,
               col_w=[2.1, 1.3, 6.9, 1.7], font=10)
        _txt(s, 0.7, 6.7, 11.9, 0.4,
             "Deskripsi dipotong agar muat di slide — rincian lengkap ada di issue log.",
             size=9.5, color=MUTED)
        hal += 1
    return hal


def _slide_breakdown(prs, nama, d, hal):
    """Sebaran status pada tracker satu proyek."""
    bd = {k: v for k, v in (d.get("breakdown") or {}).items() if str(k).strip()}
    if not bd:
        return hal
    bd = dict(sorted(bd.items(), key=lambda kv: -kv[1])[:8])

    s = _slide_kosong(prs)
    _header(s, f"Detail — {nama}", "Sebaran status pekerjaan pada tracker", hal)

    cd = CategoryChartData()
    cd.categories = [_potong(k, 22) for k in bd.keys()]
    cd.add_series("Jumlah", list(bd.values()))
    gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.7),
                            Inches(1.8), Inches(8.2), Inches(4.3), cd)
    chart = gf.chart
    chart.has_legend = False
    _rapikan_chart(chart)
    plot = chart.plots[0]
    plot.gap_width = 80
    plot.has_data_labels = True
    plot.data_labels.font.size = Pt(11)
    plot.data_labels.font.bold = True
    plot.data_labels.font.color.rgb = INK
    plot.series[0].format.fill.solid()
    plot.series[0].format.fill.fore_color.rgb = _hex_to_rgb(d["cfg"].get("warna", "#3498db"))

    _kpi(s, 9.3, 1.9, 3.3, f"{d['persen']:.0f}%", "progres tercatat")
    _kpi(s, 9.3, 3.2, 3.3, f"{d['selesai']}/{d['total_aktif']}", "item selesai / total")
    _kpi(s, 9.3, 4.5, 3.3, str(d["iss_open"]), "isu terbuka",
         warna=BAD if d["iss_open"] else OK)
    return hal + 1


def build_pptx(data, judul="Progres Proyek dan Isu Terbuka", tanggal=None,
               sertakan_breakdown=True):
    """Bangun deck dan kembalikan BytesIO siap dipakai st.download_button."""
    if tanggal is None:
        now = datetime.now()
        tanggal = f"{now.day} {BULAN_ID[now.month - 1]} {now.year}"

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    _slide_cover(prs, data, judul, tanggal)
    hal = 2
    _slide_ringkasan(prs, data, hal); hal += 1
    _slide_chart_progress(prs, data, hal); hal += 1
    _slide_chart_issue(prs, data, hal); hal += 1
    hal = _slide_issue_open(prs, data, hal)

    if sertakan_breakdown:
        for nama, d in data.items():
            hal = _slide_breakdown(prs, nama, d, hal)

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


# ═══════════════════════════════════════════════════════════
# MAIN APP
# ═══════════════════════════════════════════════════════════
try:
    col_title, col_btn = st.columns([5, 1])
    with col_title:
        st.title("📊 Dashboard Overview Project")
        st.caption("Gabungan progress & issue: DSSP · KCS CIMB Batam · KDS · Nova · Uniqlo")
    with col_btn:
        st.markdown("<br>", unsafe_allow_html=True)
        if st.button("🔄 Muat Ulang", use_container_width=True):
            load_tracker.clear()
            load_issue_log.clear()
            st.rerun()

    st.markdown("---")

    # ── Load semua data ──
    data = {}
    with st.spinner("Memuat data semua project..."):
        for nama, cfg in PROJECTS.items():
            df_track = load_tracker(cfg["tracker_sheet"])
            df_issue = load_issue_log(cfg["issue_sheet"])
            persen, total_aktif, selesai, breakdown = hitung_progress(
                df_track, cfg["progress_mode"], cfg.get("extra_selesai", ()))
            iss_total, iss_open, iss_close = hitung_issue(df_issue)
            data[nama] = {
                "cfg": cfg,
                "df_track": df_track,
                "df_issue": df_issue,
                "persen": persen,
                "total_aktif": total_aktif,
                "selesai": selesai,
                "breakdown": breakdown,
                "iss_total": iss_total,
                "iss_open": iss_open,
                "iss_close": iss_close,
            }

    # ══════════════════════════════════════════════════════
    # SECTION 1: RINGKASAN EKSEKUTIF (angka besar)
    # ══════════════════════════════════════════════════════
    total_open_all = sum(d["iss_open"] for d in data.values())
    total_issue_all = sum(d["iss_total"] for d in data.values())
    avg_progress = sum(d["persen"] for d in data.values()) / len(data) if data else 0

    m1, m2, m3, m4 = st.columns(4)
    m1.metric("📁 Total Project Aktif", len(data))
    m2.metric("📈 Rata-rata Progress", f"{avg_progress:.1f}%")
    m3.metric("🐞 Total Issue", total_issue_all)
    m4.metric("🔴 Issue OPEN (perlu tindak)", total_open_all,
              delta=None if total_open_all == 0 else f"{total_open_all} butuh aksi",
              delta_color="inverse")

    st.markdown("---")

    # ══════════════════════════════════════════════════════
    # SECTION 1B: EKSPOR KE POWERPOINT
    # ══════════════════════════════════════════════════════
    with st.expander("📤 Ekspor laporan ke PowerPoint (.pptx)"):
        ex1, ex2 = st.columns([2, 1])
        with ex1:
            judul_deck = st.text_input("Judul deck", "Progres Proyek dan Isu Terbuka")
            pilih_proyek = st.multiselect(
                "Proyek yang disertakan", list(data.keys()), default=list(data.keys()),
                help="Keluarkan proyek yang belum berjalan agar grafik tidak menyesatkan.")
        with ex2:
            tgl_deck = st.text_input("Tanggal data", datetime.now().strftime("%d %B %Y"))
            ikut_detail = st.checkbox("Slide breakdown status per proyek", True)

        if not pilih_proyek:
            st.warning("Pilih minimal satu proyek.")
        elif st.button("🛠️ Buat file PPTX"):
            subset = {k: v for k, v in data.items() if k in pilih_proyek}
            with st.spinner("Menyusun slide..."):
                st.session_state["pptx_buf"] = build_pptx(
                    subset, judul=judul_deck, tanggal=tgl_deck,
                    sertakan_breakdown=ikut_detail).getvalue()
            st.success("File siap diunduh.")

        if st.session_state.get("pptx_buf"):
            st.download_button(
                "⬇️ Unduh PPTX",
                st.session_state["pptx_buf"],
                file_name=f"Progress_Issue_Report_{datetime.now():%d%m%Y}.pptx",
                mime=("application/vnd.openxmlformats-officedocument"
                      ".presentationml.presentation"),
                use_container_width=False,
            )
        st.caption(
            "ℹ️ Deck ini berisi data mentah (grafik + daftar isu), bukan laporan naratif. "
            "Analisa, konteks MoM/email, dan rekomendasi tetap perlu ditambahkan manual."
        )

    st.markdown("---")

    # ══════════════════════════════════════════════════════
    # SECTION 2: PROGRESS GABUNGAN (side-by-side dgn issue)
    # ══════════════════════════════════════════════════════
    st.subheader("🎯 Progress & Issue per Project")

    col_prog, col_iss = st.columns(2)

    # -- Grafik progress (horizontal bar) --
    with col_prog:
        st.markdown("**Persentase Penyelesaian**")
        prog_df = pd.DataFrame({
            "Project": list(data.keys()),
            "Progress": [d["persen"] for d in data.values()],
        }).sort_values("Progress")
        fig_prog = px.bar(
            prog_df, x="Progress", y="Project", orientation="h",
            text=prog_df["Progress"].map(lambda v: f"{v:.0f}%"),
            color="Project",
            color_discrete_map={n: d["cfg"]["warna"] for n, d in data.items()},
        )
        fig_prog.update_traces(textposition="outside", cliponaxis=False)
        fig_prog.update_layout(
            showlegend=False, xaxis_title="% Selesai", yaxis_title="",
            xaxis_range=[0, 110], margin=dict(t=10, b=0), height=300,
        )
        st.plotly_chart(fig_prog, use_container_width=True)

    # -- Grafik issue open vs close (stacked) --
    with col_iss:
        st.markdown("**Status Issue (Open vs Close)**")
        iss_rows = []
        for n, d in data.items():
            iss_rows.append({"Project": n, "Kategori": "Open", "Jumlah": d["iss_open"]})
            iss_rows.append({"Project": n, "Kategori": "Close", "Jumlah": d["iss_close"]})
        iss_df = pd.DataFrame(iss_rows)
        fig_iss = px.bar(
            iss_df, x="Jumlah", y="Project", color="Kategori", orientation="h",
            barmode="stack",
            color_discrete_map={"Open": "#e74c3c", "Close": "#2ecc71"},
        )
        fig_iss.update_layout(
            yaxis_title="", xaxis_title="Jumlah Issue", legend_title="",
            margin=dict(t=10, b=0), height=300,
            legend=dict(orientation="h", yanchor="bottom", y=1.02, xanchor="right", x=1),
        )
        st.plotly_chart(fig_iss, use_container_width=True)

    st.caption(
        "ℹ️ Catatan metode: progress DSSP/KDS/Nova = rasio status *selesai* terhadap total baris aktif; "
        "KCS Batam = rata-rata kolom *Progres* numerik (0–1). "
        "Untuk DSSP, status *Under Review* dihitung sebagai selesai. "
        "Untuk KDS, status *Invoiced* dihitung sebagai pekerjaan lapangan selesai. "
        "Unit kerja tiap project berbeda (task vs site vs tiket), jadi bandingkan tren, bukan angka absolut."
    )

    st.markdown("---")

    # ══════════════════════════════════════════════════════
    # SECTION 3: KARTU RINGKAS TIAP PROJECT
    # ══════════════════════════════════════════════════════
    st.subheader("🗂️ Ringkasan per Project")
    cards = st.columns(len(data))
    for card, (nama, d) in zip(cards, data.items()):
        with card:
            st.markdown(f"#### {nama}")
            st.progress(min(d["persen"] / 100, 1.0))
            st.markdown(
                f"**{d['persen']:.0f}%** selesai &nbsp;·&nbsp; "
                f"{d['selesai']}/{d['total_aktif']} item",
                unsafe_allow_html=True,
            )
            # Tampilkan state 'Under Review' (hampir selesai) bila ada
            review_n = sum(v for k, v in d["breakdown"].items()
                           if 'review' in str(k).lower())
            if review_n > 0:
                st.caption(f"🟢 termasuk {review_n} item Under Review (dihitung selesai)")
            if d["iss_open"] > 0:
                st.error(f"🔴 {d['iss_open']} issue OPEN")
            elif d["iss_total"] > 0:
                st.success(f"✅ {d['iss_total']} issue, semua close")
            else:
                st.info("Tidak ada issue tercatat")

    st.markdown("---")

    # ══════════════════════════════════════════════════════
    # SECTION 4: DAFTAR ISSUE OPEN (actionable, digabung)
    # ══════════════════════════════════════════════════════
    st.subheader("🔴 Daftar Issue OPEN — Perlu Tindak Lanjut")

    open_rows = []
    for nama, d in data.items():
        df_iss = d["df_issue"]
        if df_iss.empty:
            continue
        col_status = _cari_kolom(df_iss.columns, 'status')
        col_desc = _cari_kolom(df_iss.columns, 'deskripsi', 'description', 'issue')
        col_pic = _cari_kolom(df_iss.columns, 'pic', 'penanggung')
        # 'remark (close)' harus menang; 'need support' TIDAK boleh dianggap remark
        col_remark = _cari_kolom(df_iss.columns, 'remark (close)', 'remark', 'keterangan', 'catatan')
        col_support = _cari_kolom(df_iss.columns, 'need support', 'support')
        col_tgl = _cari_kolom(df_iss.columns, 'tanggal', 'date')
        if not col_status:
            continue
        mask_open = df_iss[col_status].astype(str).str.strip().str.lower().str.contains('open', na=False)
        for _, r in df_iss[mask_open].iterrows():
            open_rows.append({
                "Project": nama,
                "Tanggal": str(r[col_tgl]) if col_tgl else "",
                "Deskripsi": str(r[col_desc]) if col_desc else "",
                "PIC": str(r[col_pic]) if col_pic else "",
                "Need Support": str(r[col_support]) if col_support else "",
                "Remark": str(r[col_remark]) if col_remark else "",
            })

    if open_rows:
        df_open = pd.DataFrame(open_rows)
        st.error(f"Terdapat {len(df_open)} issue OPEN di seluruh project.")
        st.dataframe(df_open, use_container_width=True, hide_index=True)
    else:
        st.success("✅ Tidak ada issue OPEN. Semua issue tercatat sudah close.")

    st.markdown("---")

    # ══════════════════════════════════════════════════════
    # SECTION 5: DETAIL PER PROJECT (expander)
    # ══════════════════════════════════════════════════════
    st.subheader("🔍 Detail Data per Project")
    for nama, d in data.items():
        with st.expander(f"{nama} — {d['persen']:.0f}% selesai · {d['iss_open']} issue open"):
            punya_margin = ('KDS' in nama) or ('Nova' in nama)
            punya_cost_health = ('Uniqlo' in nama)
            if punya_margin:
                tab_t, tab_i, tab_s, tab_m = st.tabs(
                    ["📋 Tracker", "🐞 Issue Log", "📊 Breakdown Status", "💰 Margin Profit"])
            elif punya_cost_health:
                tab_t, tab_i, tab_s, tab_c = st.tabs(
                    ["📋 Tracker", "🐞 Issue Log", "📊 Breakdown Status", "💵 Kesehatan Cost"])
            else:
                tab_t, tab_i, tab_s = st.tabs(
                    ["📋 Tracker", "🐞 Issue Log", "📊 Breakdown Status"])

            with tab_t:
                if not d["df_track"].empty:
                    st.dataframe(d["df_track"], use_container_width=True, height=360)
                else:
                    st.warning("Data tracker kosong.")

            with tab_i:
                if not d["df_issue"].empty:
                    st.dataframe(d["df_issue"], use_container_width=True, height=300)
                else:
                    st.info("Belum ada issue tercatat untuk project ini.")

            with tab_s:
                if d["breakdown"]:
                    bd = pd.DataFrame(
                        {"Status": list(d["breakdown"].keys()),
                         "Jumlah": list(d["breakdown"].values())}
                    )
                    color_map = {
                        'Done': '#2ecc71', 'Selesai': '#2ecc71', 'Invoiced': '#2ecc71',
                        'On Progress': '#3498db', 'In Progress': '#3498db',
                        'Under Review': '#f39c12',
                        'TBD': '#95a5a6', 'NYS': '#95a5a6', 'Not yet': '#95a5a6',
                        'Pending': '#e74c3c', 'Unknown': '#bdc3c7',
                    }
                    fig_bd = px.bar(bd, x="Status", y="Jumlah", color="Status",
                                    color_discrete_map=color_map, text="Jumlah")
                    fig_bd.update_traces(textposition="outside", cliponaxis=False)
                    fig_bd.update_layout(showlegend=False, margin=dict(t=10, b=0), height=320)
                    st.plotly_chart(fig_bd, use_container_width=True)
                else:
                    st.info("Tidak ada data status untuk ditampilkan.")

            if punya_margin:
                with tab_m:
                    rincian, tot_rev, tot_cost, tot_margin = compute_margin(d["df_track"], nama)
                    if rincian is None or rincian.empty:
                        st.info("Data revenue/cost tidak tersedia untuk menghitung margin.")
                    else:
                        if 'Nova' in nama:
                            st.caption("ℹ️ Nova: Revenue = PO, Cost = Actual Cost, "
                                       "Margin = PO − Actual Cost.")
                        else:
                            st.caption("ℹ️ KDS: Revenue, Cost, dan Margin dari kolom tracker.")

                        pct = (tot_margin / tot_rev * 100) if tot_rev else 0
                        mk1, mk2, mk3, mk4 = st.columns(4)
                        mk1.metric("💰 Total Revenue", f"Rp {tot_rev:,.0f}")
                        mk2.metric("📉 Total Cost", f"Rp {tot_cost:,.0f}")
                        mk3.metric("📈 Total Margin", f"Rp {tot_margin:,.0f}")
                        mk4.metric("📊 Margin %", f"{pct:.1f}%")

                        # Peringatan bila ada lokasi rugi (margin negatif)
                        rugi = rincian[rincian['Margin'] < 0]
                        if not rugi.empty:
                            st.warning(f"⚠️ {len(rugi)} lokasi margin negatif (rugi). "
                                       "Cek tabel di bawah.")

                        # Grafik margin per lokasi (hijau = untung, merah = rugi)
                        viz = rincian.sort_values('Margin').copy()
                        viz['Warna'] = viz['Margin'].apply(
                            lambda v: 'Rugi' if v < 0 else 'Untung')
                        fig_m = px.bar(
                            viz, x='Margin', y='Lokasi', orientation='h', color='Warna',
                            color_discrete_map={'Untung': '#2ecc71', 'Rugi': '#e74c3c'},
                            height=max(320, len(viz) * 28),
                        )
                        fig_m.update_layout(
                            showlegend=False, yaxis_title="", xaxis_title="Margin (Rp)",
                            margin=dict(t=10, b=0))
                        st.plotly_chart(fig_m, use_container_width=True)

                        # Tabel rincian
                        tampil = rincian.copy()
                        for c in ['Revenue', 'Cost', 'Margin']:
                            tampil[c] = tampil[c].map(lambda v: f"Rp {v:,.0f}")
                        tampil['Margin %'] = rincian['Margin %'].map(lambda v: f"{v:.1f}%")
                        st.dataframe(tampil, use_container_width=True, hide_index=True, height=300)

            if punya_cost_health:
                with tab_c:
                    ch = compute_cost_health(d["df_track"], nama)
                    if ch is None:
                        st.info("Kolom Plan Cost / Actual Cost tidak ditemukan di tracker.")
                    else:
                        st.caption(
                            "ℹ️ Kesehatan cost = Plan vs Actual. **Headroom pool** = "
                            "Σ Plan − Σ Actual dari lokasi yg SUDAH realisasi. Selama headroom "
                            "positif, lokasi overcost masih bisa disubsidi silang oleh lokasi "
                            "undercost. Lokasi tanpa Actual Cost belum dihitung (lihat sub-tabel)."
                        )

                        real = ch['realized']
                        if real.empty:
                            st.warning(
                                "⏳ Belum ada lokasi dengan Actual Cost terisi. "
                                f"Semua {len(ch['unrealized'])} lokasi masih *belum realisasi*, "
                                "jadi kesehatan cost belum bisa dihitung. Panel ini akan aktif "
                                "begitu kolom **Actual Cost** mulai diisi."
                            )
                        else:
                            c1, c2, c3, c4 = st.columns(4)
                            c1.metric("💰 Plan (realisasi)", f"Rp {ch['plan_real']:,.0f}")
                            c2.metric("💸 Actual (realisasi)", f"Rp {ch['actual_real']:,.0f}")
                            c3.metric(
                                "🏦 Headroom Pool",
                                f"Rp {ch['headroom']:,.0f}",
                                delta=("sisa anggaran" if ch['headroom'] >= 0
                                       else "pool defisit"),
                                delta_color=("normal" if ch['headroom'] >= 0 else "inverse"),
                            )
                            c4.metric("🔴 Lokasi Overcost", ch['n_over'])

                            # Inti keputusan subsidi silang
                            if ch['n_over'] == 0:
                                st.success(
                                    "✅ Tidak ada lokasi overcost. Semua realisasi di bawah/tepat plan."
                                )
                            elif ch['headroom'] >= 0:
                                st.success(
                                    f"✅ {ch['n_over']} lokasi overcost (total kelebihan "
                                    f"Rp {ch['over_total']:,.0f}), TAPI masih tertutup headroom "
                                    f"Rp {ch['headroom']:,.0f}. Subsidi silang layak — "
                                    "pool anggaran agregat masih surplus."
                                )
                            else:
                                st.error(
                                    f"⚠️ {ch['n_over']} lokasi overcost (total kelebihan "
                                    f"Rp {ch['over_total']:,.0f}) dan headroom sudah defisit "
                                    f"Rp {ch['headroom']:,.0f}. Overcost TIDAK lagi tertutup "
                                    "oleh lokasi undercost — perlu tambahan anggaran / evaluasi."
                                )

                            # Grafik variance per lokasi (merah=overcost, hijau=undercost)
                            viz = real.sort_values('Variance').copy()
                            fig_c = px.bar(
                                viz, x='Variance', y='Lokasi', orientation='h',
                                color='Kondisi',
                                color_discrete_map={
                                    'Overcost': '#e74c3c',
                                    'Undercost': '#2ecc71',
                                    'Tepat': '#95a5a6',
                                },
                                height=max(320, len(viz) * 28),
                                labels={'Variance': 'Selisih Actual − Plan (Rp)'},
                            )
                            fig_c.update_layout(
                                yaxis_title="", xaxis_title="Selisih thd Plan (Rp) · + = overcost",
                                legend_title="", margin=dict(t=10, b=0),
                                legend=dict(orientation="h", yanchor="bottom",
                                            y=1.02, xanchor="right", x=1),
                            )
                            fig_c.add_vline(x=0, line_width=1, line_color="#333")
                            st.plotly_chart(fig_c, use_container_width=True)

                            # Tabel lokasi terealisasi
                            tampil_r = real[['Lokasi', 'Kota', 'Plan', 'Actual',
                                             'Variance', 'Variance %', 'Kondisi']].copy()
                            for c in ['Plan', 'Actual', 'Variance']:
                                tampil_r[c] = tampil_r[c].map(lambda v: f"Rp {v:,.0f}")
                            tampil_r['Variance %'] = real['Variance %'].map(
                                lambda v: f"{v:+.1f}%")
                            st.markdown("**Lokasi sudah realisasi**")
                            st.dataframe(tampil_r, use_container_width=True,
                                         hide_index=True, height=280)

                        # Sub-tabel: lokasi belum realisasi (selalu tampil bila ada)
                        unreal = ch['unrealized']
                        if not unreal.empty:
                            st.markdown(f"**⏳ Belum realisasi ({len(unreal)} lokasi)** "
                                        "— Actual Cost belum diisi, tidak masuk hitungan headroom")
                            tampil_u = unreal[['Lokasi', 'Kota', 'Plan']].copy()
                            tampil_u['Plan'] = tampil_u['Plan'].map(lambda v: f"Rp {v:,.0f}")
                            tampil_u['Actual'] = "— belum diisi —"
                            st.dataframe(tampil_u, use_container_width=True,
                                         hide_index=True, height=240)

except Exception as e:
    st.error(f"Terjadi kesalahan sistem: {e}")
    st.exception(e)
